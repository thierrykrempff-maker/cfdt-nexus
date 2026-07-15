"""Local, non-destructive document importer without OCR or network access."""

from __future__ import annotations

import argparse
import importlib.metadata
import importlib.util
import json
import os
import re
import shutil
import subprocess
import time
import uuid
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable, Iterable

from automation.cse_memory.audit_cse_corpus import detect_family, detect_years, sha256_file
from automation.cse_memory.document_models import DocumentRecord, ExtractionResult, SCHEMA_VERSION


DIRECT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".txt"}
CONVERTER_REQUIRED = {".doc", ".ppt", ".xls", ".rtf", ".msg"}
UNSUPPORTED = {
    ".zip", ".7z", ".rar", ".jpg", ".jpeg", ".png", ".gif", ".bmp",
    ".tif", ".tiff", ".webp", ".db", ".lnk", ".partial",
}
EXTRACTOR_PACKAGES = {
    ".pdf": ("pypdf", "pypdf"),
    ".docx": ("docx", "python-docx"),
    ".pptx": ("pptx", "python-pptx"),
    ".xlsx": ("openpyxl", "openpyxl"),
}
MAX_XLSX_CELLS = 100_000
ABSOLUTE_PATH_PATTERN = re.compile(
    r"(?i)(?<![\w])(?:[a-z]:[\\/][^\r\n\t<>\"|?*]+|/(?:home|users)/[^\r\n\t<>\"|?*]+)"
)


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def stable_document_id(relative_path: str) -> str:
    normalized = Path(relative_path).as_posix()
    return str(uuid.uuid5(uuid.NAMESPACE_URL, f"cfdt-nexus:cse:{normalized}"))


def dependency_version(distribution: str) -> str | None:
    try:
        return importlib.metadata.version(distribution)
    except importlib.metadata.PackageNotFoundError:
        return None


def extraction_capabilities() -> dict[str, dict[str, Any]]:
    formats = {
        "PDF": (".pdf", "direct"), "DOCX": (".docx", "direct"),
        "DOC": (".doc", "converter_required"), "PPTX": (".pptx", "direct"),
        "PPT": (".ppt", "converter_required"), "XLSX": (".xlsx", "direct"),
        "XLS": (".xls", "converter_required"), "TXT": (".txt", "direct"),
        "RTF": (".rtf", "converter_required"), "MSG": (".msg", "converter_required"),
        "ZIP": (".zip", "unsupported"), "images": (".jpg", "unsupported"),
    }
    result: dict[str, dict[str, Any]] = {}
    for label, (extension, planned) in formats.items():
        module, distribution = EXTRACTOR_PACKAGES.get(extension, (None, None))
        available = True if extension == ".txt" else bool(module and importlib.util.find_spec(module))
        if planned == "direct" and not available:
            status = "not_available"
        else:
            status = planned
        result[label] = {
            "status": status,
            "module": module,
            "version": dependency_version(distribution) if distribution else None,
            "notes": "No OCR" if label in {"PDF", "images"} else None,
        }
    return result


def _clean_text(parts: Iterable[str]) -> str:
    return "\n".join(part.strip() for part in parts if part and part.strip())


def redact_absolute_paths(value: Any) -> tuple[Any, bool]:
    if isinstance(value, str):
        cleaned, count = ABSOLUTE_PATH_PATTERN.subn("[ABSOLUTE_PATH_REDACTED]", value)
        return cleaned, count > 0
    if isinstance(value, list):
        changed = False
        cleaned = []
        for item in value:
            safe, item_changed = redact_absolute_paths(item)
            cleaned.append(safe)
            changed = changed or item_changed
        return cleaned, changed
    if isinstance(value, dict):
        changed = False
        cleaned = {}
        for key, item in value.items():
            safe, item_changed = redact_absolute_paths(item)
            cleaned[key] = safe
            changed = changed or item_changed
        return cleaned, changed
    return value, False


def extract_pdf(path: Path) -> ExtractionResult:
    from pypdf import PdfReader

    reader = PdfReader(path)
    parts: list[str] = []
    warnings: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if not text.strip():
            warnings.append(f"page_without_text:{index}")
        parts.append(f"--- PAGE {index} ---\n{text.strip()}")
    metadata = {str(key).lstrip("/"): str(value) for key, value in (reader.metadata or {}).items()}
    return ExtractionResult(
        text_content="\n\n".join(parts),
        status="extracted_with_warnings" if warnings else "extracted",
        extractor_name="pypdf",
        extractor_method="page.extract_text (no OCR)",
        page_count=len(reader.pages),
        metadata=metadata,
        warnings=warnings,
    )


def extract_docx(path: Path) -> ExtractionResult:
    from docx import Document

    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table_index, table in enumerate(document.tables, start=1):
        parts.append(f"--- TABLE {table_index} ---")
        parts.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
    properties = document.core_properties
    metadata = {
        key: str(value) for key, value in {
            "title": properties.title, "subject": properties.subject,
            "category": properties.category, "created": properties.created,
            "modified": properties.modified,
        }.items() if value is not None
    }
    return ExtractionResult(
        text_content=_clean_text(parts), extractor_name="python-docx",
        extractor_method="paragraphs and tables", paragraph_count=len(document.paragraphs),
        metadata={**metadata, "table_count": len(document.tables)},
    )


def extract_pptx(path: Path) -> ExtractionResult:
    from pptx import Presentation

    presentation = Presentation(path)
    slides: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts = [f"--- SLIDE {slide_index} ---"]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.extend(paragraph.text for paragraph in shape.text_frame.paragraphs)
            if getattr(shape, "has_table", False):
                parts.extend("\t".join(cell.text for cell in row.cells) for row in shape.table.rows)
        slides.append(_clean_text(parts))
    properties = presentation.core_properties
    return ExtractionResult(
        text_content="\n\n".join(slides), extractor_name="python-pptx",
        extractor_method="shape text and tables", slide_count=len(presentation.slides),
        metadata={"title": properties.title or "", "subject": properties.subject or ""},
    )


def extract_xlsx(path: Path) -> ExtractionResult:
    import openpyxl

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True, keep_links=False)
    parts: list[str] = []
    warnings: list[str] = []
    cell_count = 0
    try:
        for worksheet in workbook.worksheets:
            parts.append(f"--- SHEET {worksheet.title} ---")
            if getattr(getattr(worksheet, "protection", None), "sheet", False):
                warnings.append(f"protected_sheet:{worksheet.title}")
            for row in worksheet.iter_rows():
                values = [str(cell.value) if cell.value is not None else "" for cell in row]
                cell_count += len(row)
                if any(values):
                    parts.append("\t".join(values))
                if cell_count >= MAX_XLSX_CELLS:
                    warnings.append(f"cell_limit_reached:{MAX_XLSX_CELLS}")
                    break
            if cell_count >= MAX_XLSX_CELLS:
                break
        return ExtractionResult(
            text_content=_clean_text(parts),
            status="extracted_with_warnings" if warnings else "extracted",
            extractor_name="openpyxl", extractor_method="read_only data_only cell values",
            sheet_count=len(workbook.sheetnames),
            metadata={"sheet_names": list(workbook.sheetnames), "cells_examined": cell_count},
            warnings=warnings,
        )
    finally:
        workbook.close()


def extract_txt(path: Path) -> ExtractionResult:
    data = path.read_bytes()
    failures: list[str] = []
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            text = data.decode(encoding)
            warnings = [] if encoding in {"utf-8", "utf-8-sig"} else [f"fallback_encoding:{encoding}"]
            return ExtractionResult(
                text_content=text, status="extracted_with_warnings" if warnings else "extracted",
                extractor_name="python", extractor_method=f"decode:{encoding}",
                metadata={"encoding": encoding}, warnings=warnings,
            )
        except UnicodeDecodeError:
            failures.append(encoding)
    raise UnicodeError(f"No supported encoding matched: {', '.join(failures)}")


EXTRACTORS: dict[str, Callable[[Path], ExtractionResult]] = {
    ".pdf": extract_pdf, ".docx": extract_docx, ".pptx": extract_pptx,
    ".xlsx": extract_xlsx, ".txt": extract_txt,
}


def _unsupported_result(extension: str) -> ExtractionResult:
    if extension in CONVERTER_REQUIRED:
        return ExtractionResult(
            status="converter_required", extractor_name="none",
            extractor_method="safe local converter unavailable",
            warnings=[f"converter_required:{extension}"],
        )
    reason = "archive_not_expanded" if extension in {".zip", ".7z", ".rar"} else "format_not_supported"
    if extension in {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tif", ".tiff", ".webp"}:
        reason = "ocr_forbidden"
    if extension == ".lnk":
        reason = "windows_shortcut_not_opened"
    return ExtractionResult(
        status="unsupported", extractor_name="none", extractor_method="not_opened",
        warnings=[f"{reason}:{extension or '[no_extension]'}"],
    )


def build_record(path: Path, source_root: Path) -> DocumentRecord:
    relative = path.relative_to(source_root).as_posix()
    extension = path.suffix.casefold()
    stat = path.stat()
    digest = sha256_file(path)
    years = detect_years(relative)
    base = dict(
        document_id=stable_document_id(relative), source_relative_path=relative,
        source_filename=path.name, source_extension=extension,
        source_size_bytes=stat.st_size, source_sha256=digest,
        source_modified_at=datetime.fromtimestamp(stat.st_mtime, timezone.utc).isoformat(),
        detected_year=years[-1] if years else None, detected_family=detect_family(relative),
        imported_at=utc_now(), schema_version=SCHEMA_VERSION,
    )
    if stat.st_size == 0:
        result = ExtractionResult(status="empty", extractor_name="none", extractor_method="size_check")
        error_code = error_message = None
    else:
        extractor = EXTRACTORS.get(extension)
        if extractor is None:
            result = _unsupported_result(extension)
            error_code = error_message = None
        elif extension in EXTRACTOR_PACKAGES and importlib.util.find_spec(EXTRACTOR_PACKAGES[extension][0]) is None:
            result = ExtractionResult(
                status="converter_required", extractor_name="none",
                extractor_method="dependency_unavailable", warnings=["extractor_dependency_unavailable"],
            )
            error_code = error_message = None
        else:
            try:
                result = extractor(path)
                error_code = error_message = None
            except (OSError, ValueError, KeyError, TypeError, UnicodeError) as error:
                result = ExtractionResult(
                    status="unreadable",
                    extractor_name=getattr(extractor, "__name__", type(extractor).__name__),
                    extractor_method="direct",
                )
                error_code = type(error).__name__
                error_message = str(error)[:500]
            except Exception as error:  # isolate malformed third-party parser failures
                result = ExtractionResult(
                    status="failed",
                    extractor_name=getattr(extractor, "__name__", type(extractor).__name__),
                    extractor_method="direct",
                )
                error_code = type(error).__name__
                error_message = str(error)[:500]
    safe_text, text_redacted = redact_absolute_paths(result.text_content)
    safe_metadata, metadata_redacted = redact_absolute_paths(result.metadata)
    safe_error, error_redacted = redact_absolute_paths(error_message)
    if text_redacted or metadata_redacted or error_redacted:
        result.warnings.append("absolute_path_redacted")
        if result.status == "extracted":
            result.status = "extracted_with_warnings"
    return DocumentRecord(
        **base, extractor_name=result.extractor_name, extractor_method=result.extractor_method,
        extraction_status=result.status, extraction_error_code=error_code,
        extraction_error_message=safe_error, text_content=safe_text,
        text_length=len(safe_text), page_count=result.page_count,
        sheet_count=result.sheet_count, slide_count=result.slide_count,
        paragraph_count=result.paragraph_count, technical_metadata=safe_metadata,
        warnings=result.warnings,
    )


def validate_paths(source: Path, output: Path) -> tuple[Path, Path]:
    source = source.resolve()
    output = output.resolve()
    if not source.is_dir():
        raise FileNotFoundError(f"Source directory does not exist: {source}")
    if output == source or source in output.parents:
        raise ValueError("Output must not be located inside RAW_DOCUMENTS/source")
    return source, output


def output_ignore_warning(source: Path, output: Path) -> str | None:
    """Return a relative, non-sensitive warning when Git protection is not confirmed."""
    repository = next((parent for parent in (source, *source.parents) if (parent / ".git").exists()), None)
    if repository is None:
        return "gitignore_protection_not_verifiable:no_repository"
    try:
        relative = output.relative_to(repository).as_posix()
    except ValueError:
        return "gitignore_protection_not_verifiable:output_outside_repository"
    git = shutil.which("git")
    if git:
        result = subprocess.run(
            [git, "check-ignore", "-q", "--", relative], cwd=repository,
            check=False, capture_output=True, timeout=10,
        )
        if result.returncode == 0:
            return None
    ignore_file = repository / ".gitignore"
    try:
        patterns = {
            line.strip().rstrip("/") for line in ignore_file.read_text(encoding="utf-8").splitlines()
            if line.strip() and not line.lstrip().startswith("#") and not any(char in line for char in "*?[]!")
        }
        if any(relative == pattern or relative.startswith(f"{pattern}/") for pattern in patterns):
            return None
    except OSError:
        pass
    return f"output_not_confirmed_gitignored:{relative}"


def discover_files(source: Path, extensions: set[str] | None = None, subfolder: str | None = None) -> list[Path]:
    scan_root = source / subfolder if subfolder else source
    if not scan_root.is_dir() or source not in (scan_root.resolve(), *scan_root.resolve().parents):
        raise ValueError("Subfolder must exist inside source")
    files: list[Path] = []
    for current, directories, filenames in os.walk(scan_root, followlinks=False):
        directories[:] = [name for name in directories if not (Path(current) / name).is_symlink()]
        for filename in filenames:
            path = Path(current) / filename
            if path.is_symlink() or path.suffix.casefold() == ".lnk":
                if path.suffix.casefold() == ".lnk" and (not extensions or ".lnk" in extensions):
                    files.append(path)
                continue
            if not extensions or path.suffix.casefold() in extensions:
                files.append(path)
    return sorted(files, key=lambda item: item.relative_to(source).as_posix().casefold())


def _write_json(path: Path, value: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _summary_markdown(manifest: dict[str, Any]) -> str:
    counts = "\n".join(f"- {key}: {value}" for key, value in manifest["status_counts"].items()) or "- Aucun"
    return f"""# Import documentaire local — LOT 1A

- Mode : {manifest['mode']}
- Fichiers examinés : {manifest['examined_count']}
- Fichiers traités : {manifest['processed_count']}
- Fichiers repris sans retraitement : {manifest['resumed_count']}
- Volume textuel extrait : {manifest['text_length_total']} caractères
- Durée : {manifest['duration_seconds']:.3f} secondes

## Statuts

{counts}

Extraction locale uniquement, sans OCR, réseau, IA externe, chunks ou indexation sémantique.
"""


def run_import(
    source: Path | str,
    output: Path | str,
    *, mode: str = "dry-run", extensions: set[str] | None = None,
    subfolder: str | None = None, limit: int | None = None, force: bool = False,
    extension_limits: dict[str, int] | None = None,
) -> dict[str, Any]:
    source, output = validate_paths(Path(source), Path(output))
    if mode not in {"dry-run", "import"}:
        raise ValueError("Mode must be dry-run or import")
    normalized_extensions = {value.casefold() if value.startswith(".") else f".{value.casefold()}" for value in extensions or set()}
    files = discover_files(source, normalized_extensions or None, subfolder)
    if extension_limits:
        normalized_limits = {
            key.casefold() if key.startswith(".") else f".{key.casefold()}": max(0, value)
            for key, value in extension_limits.items()
        }
        selected: list[Path] = []
        selected_counts: Counter[str] = Counter()
        for path in files:
            extension = path.suffix.casefold()
            if extension in normalized_limits and selected_counts[extension] < normalized_limits[extension]:
                selected.append(path)
                selected_counts[extension] += 1
        files = selected
    if limit is not None:
        files = files[: max(0, limit)]
    started = time.monotonic()
    security_warnings = [warning] if (warning := output_ignore_warning(source, output)) else []
    status_counts: Counter[str] = Counter()
    manifest_documents: list[dict[str, Any]] = []
    errors: list[dict[str, Any]] = []
    resumed = 0
    processed = 0
    text_total = 0

    for path in files:
        relative = path.relative_to(source).as_posix()
        document_id = stable_document_id(relative)
        document_path = output / "documents" / f"{document_id}.json"
        if mode == "import" and document_path.is_file() and not force:
            try:
                previous = json.loads(document_path.read_text(encoding="utf-8"))
                if previous.get("source_sha256") == sha256_file(path):
                    resumed += 1
                    status = previous.get("extraction_status", "failed")
                    status_counts[status] += 1
                    text_total += int(previous.get("text_length", 0))
                    manifest_documents.append({
                        "document_id": document_id, "source_relative_path": relative,
                        "status": status, "resumed": True,
                    })
                    continue
            except (OSError, ValueError, TypeError):
                pass
        try:
            if mode == "dry-run":
                extension = path.suffix.casefold()
                status = "extractable" if extension in DIRECT_EXTENSIONS else _unsupported_result(extension).status
                manifest_documents.append({
                    "document_id": document_id, "source_relative_path": relative,
                    "source_extension": extension, "planned_status": status,
                })
                status_counts[status] += 1
                continue
            record = build_record(path, source)
            payload = record.to_dict()
            _write_json(document_path, payload)
            processed += 1
            text_total += record.text_length
            status_counts[record.extraction_status] += 1
            manifest_documents.append({
                "document_id": record.document_id,
                "source_relative_path": record.source_relative_path,
                "status": record.extraction_status, "text_length": record.text_length,
                "resumed": False,
            })
            if record.extraction_error_code:
                errors.append({
                    "document_id": record.document_id,
                    "source_relative_path": record.source_relative_path,
                    "error_code": record.extraction_error_code,
                    "error_message": record.extraction_error_message,
                })
        except Exception as error:  # one file must never abort the batch
            status_counts["failed"] += 1
            errors.append({
                "document_id": document_id, "source_relative_path": relative,
                "error_code": type(error).__name__, "error_message": str(error)[:500],
            })

    manifest = {
        "schema_version": SCHEMA_VERSION, "mode": mode, "generated_at": utc_now(),
        "examined_count": len(files), "processed_count": processed,
        "resumed_count": resumed, "status_counts": dict(sorted(status_counts.items())),
        "text_length_total": text_total,
        "duration_seconds": round(time.monotonic() - started, 3),
        "documents": manifest_documents,
        "capabilities": extraction_capabilities(),
        "security_warnings": security_warnings,
    }
    if mode == "import":
        _write_json(output / "manifests" / "import_manifest.json", manifest)
        (output / "manifests").mkdir(parents=True, exist_ok=True)
        (output / "manifests" / "import_summary.md").write_text(_summary_markdown(manifest), encoding="utf-8")
        _write_json(output / "logs" / "import_errors.json", {"errors": errors})
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--source", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--mode", choices=("dry-run", "import"), default="dry-run")
    parser.add_argument("--extension", action="append", dest="extensions")
    parser.add_argument("--subfolder")
    parser.add_argument("--limit", type=int)
    parser.add_argument(
        "--limit-per-extension", action="append", default=[], metavar="EXT=COUNT",
        help="Repeatable balanced sample limit, for example pdf=5",
    )
    parser.add_argument("--force", action="store_true")
    args = parser.parse_args()
    extension_limits: dict[str, int] = {}
    for value in args.limit_per_extension:
        try:
            extension, count = value.rsplit("=", 1)
            extension_limits[extension] = int(count)
        except (ValueError, TypeError) as error:
            parser.error(f"Invalid --limit-per-extension value {value!r}: {error}")
    manifest = run_import(
        args.source, args.output, mode=args.mode,
        extensions=set(args.extensions or []), subfolder=args.subfolder,
        limit=args.limit, force=args.force, extension_limits=extension_limits or None,
    )
    summary = {key: manifest[key] for key in (
        "mode", "examined_count", "processed_count", "resumed_count",
        "status_counts", "text_length_total", "duration_seconds",
    )}
    if args.mode == "import":
        summary["reports"] = [
            str(args.output / "manifests" / "import_manifest.json"),
            str(args.output / "manifests" / "import_summary.md"),
            str(args.output / "logs" / "import_errors.json"),
        ]
    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
