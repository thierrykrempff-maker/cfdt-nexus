import json
import socket
import tempfile
import unittest
from pathlib import Path
from unittest import mock

from automation.cse_memory import document_importer as importer


class DocumentImporterTests(unittest.TestCase):
    def setUp(self):
        self.temporary = tempfile.TemporaryDirectory()
        self.base = Path(self.temporary.name)
        self.source = self.base / "RAW_DOCUMENTS"
        self.output = self.base / "PROCESSED" / "LOT_1A"
        self.source.mkdir()

    def tearDown(self):
        self.temporary.cleanup()

    def write(self, relative: str, content: bytes = b"synthetic content") -> Path:
        path = self.source / relative
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(content)
        return path

    def test_document_id_and_sha256_are_stable_and_paths_relative(self):
        path = self.write("PV_CSE/2024/example.txt")
        first = importer.build_record(path, self.source)
        second = importer.build_record(path, self.source)
        self.assertEqual(first.document_id, second.document_id)
        self.assertEqual(first.source_sha256, second.source_sha256)
        self.assertEqual(first.source_relative_path, "PV_CSE/2024/example.txt")
        self.assertNotIn(str(self.base), json.dumps(first.to_dict()))

    def test_txt_extraction_and_fallback_encoding(self):
        path = self.write("notes.txt", "texte synthétique".encode("cp1252"))
        record = importer.build_record(path, self.source)
        self.assertEqual(record.extraction_status, "extracted_with_warnings")
        self.assertIn("texte synthétique", record.text_content)

    def test_absolute_paths_in_extracted_content_are_redacted(self):
        path = self.write("path.txt", b"Synthetic C:\\Users\\Example\\secret.txt value")
        record = importer.build_record(path, self.source)
        self.assertNotIn("C:\\Users", record.text_content)
        self.assertIn("[ABSOLUTE_PATH_REDACTED]", record.text_content)
        self.assertIn("absolute_path_redacted", record.warnings)

    def test_docx_extraction_includes_paragraphs_and_tables(self):
        from docx import Document

        path = self.source / "synthetic.docx"
        document = Document()
        document.add_paragraph("Synthetic paragraph")
        table = document.add_table(rows=1, cols=2)
        table.cell(0, 0).text = "Cell A"
        table.cell(0, 1).text = "Cell B"
        document.save(path)
        record = importer.build_record(path, self.source)
        self.assertEqual(record.extraction_status, "extracted")
        self.assertIn("Synthetic paragraph", record.text_content)
        self.assertIn("Cell A", record.text_content)
        self.assertEqual(record.technical_metadata["table_count"], 1)

    def test_pptx_extraction_and_slide_separator(self):
        from pptx import Presentation

        path = self.source / "synthetic.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Synthetic title"
        presentation.save(path)
        record = importer.build_record(path, self.source)
        self.assertEqual(record.slide_count, 1)
        self.assertIn("--- SLIDE 1 ---", record.text_content)
        self.assertIn("Synthetic title", record.text_content)

    def test_xlsx_extraction_and_sheet_count(self):
        from openpyxl import Workbook

        path = self.source / "synthetic.xlsx"
        workbook = Workbook()
        workbook.active.title = "Synthetic sheet"
        workbook.active["A1"] = "Synthetic value"
        workbook.create_sheet("Second")
        workbook.save(path)
        record = importer.build_record(path, self.source)
        self.assertEqual(record.sheet_count, 2)
        self.assertIn("Synthetic value", record.text_content)

    def test_pdf_extraction_without_ocr(self):
        from reportlab.pdfgen.canvas import Canvas

        path = self.source / "synthetic.pdf"
        canvas = Canvas(str(path))
        canvas.drawString(72, 720, "Synthetic PDF text")
        canvas.save()
        record = importer.build_record(path, self.source)
        self.assertEqual(record.page_count, 1)
        self.assertIn("Synthetic PDF text", record.text_content)
        self.assertIn("no OCR", record.extractor_method)

    def test_unsupported_converter_empty_and_shortcut(self):
        legacy = importer.build_record(self.write("legacy.doc"), self.source)
        archive = importer.build_record(self.write("archive.zip"), self.source)
        shortcut = importer.build_record(self.write("shortcut.lnk"), self.source)
        empty = importer.build_record(self.write("empty.txt", b""), self.source)
        self.assertEqual(legacy.extraction_status, "converter_required")
        self.assertEqual(archive.extraction_status, "unsupported")
        self.assertEqual(shortcut.extraction_status, "unsupported")
        self.assertIn("windows_shortcut_not_opened", shortcut.warnings[0])
        self.assertEqual(empty.extraction_status, "empty")

    def test_unreadable_file_is_recorded(self):
        path = self.write("denied.txt")
        with mock.patch.dict(importer.EXTRACTORS, {".txt": mock.Mock(side_effect=PermissionError("synthetic"))}):
            record = importer.build_record(path, self.source)
        self.assertEqual(record.extraction_status, "unreadable")
        self.assertEqual(record.extraction_error_code, "PermissionError")

    def test_import_outputs_manifest_errors_and_one_json_per_document(self):
        self.write("one.txt", b"one")
        self.write("two.zip", b"two")
        result = importer.run_import(self.source, self.output, mode="import")
        self.assertEqual(result["examined_count"], 2)
        self.assertEqual(len(list((self.output / "documents").glob("*.json"))), 2)
        self.assertTrue((self.output / "manifests" / "import_manifest.json").is_file())
        self.assertTrue((self.output / "manifests" / "import_summary.md").is_file())
        self.assertTrue((self.output / "logs" / "import_errors.json").is_file())

    def test_error_is_isolated_without_stopping_batch(self):
        bad = self.write("bad.txt", b"bad")
        self.write("good.txt", b"good")
        original = importer.build_record

        def selective(path, source):
            if path == bad:
                raise RuntimeError("synthetic failure")
            return original(path, source)

        with mock.patch.object(importer, "build_record", side_effect=selective):
            result = importer.run_import(self.source, self.output, mode="import")
        self.assertEqual(result["status_counts"]["failed"], 1)
        self.assertEqual(result["status_counts"]["extracted"], 1)

    def test_resume_and_force_reimport(self):
        self.write("resume.txt", b"stable")
        first = importer.run_import(self.source, self.output, mode="import")
        second = importer.run_import(self.source, self.output, mode="import")
        forced = importer.run_import(self.source, self.output, mode="import", force=True)
        self.assertEqual(first["processed_count"], 1)
        self.assertEqual(second["resumed_count"], 1)
        self.assertEqual(second["processed_count"], 0)
        self.assertEqual(forced["processed_count"], 1)

    def test_refuses_missing_source_and_output_inside_source(self):
        with self.assertRaises(FileNotFoundError):
            importer.run_import(self.base / "missing", self.output)
        with self.assertRaises(ValueError):
            importer.run_import(self.source, self.source / "output")

    def test_warns_when_output_is_not_confirmed_gitignored(self):
        self.write("one.txt")
        result = importer.run_import(self.source, self.output, mode="dry-run")
        self.assertTrue(result["security_warnings"])

    def test_originals_unchanged_and_no_absolute_paths_or_network(self):
        source = self.write("immutable.txt", b"immutable")
        before = (source.read_bytes(), source.stat().st_mtime_ns)
        with mock.patch.object(socket, "create_connection", side_effect=AssertionError("network forbidden")) as network:
            importer.run_import(self.source, self.output, mode="import")
        after = (source.read_bytes(), source.stat().st_mtime_ns)
        self.assertEqual(before, after)
        network.assert_not_called()
        for output_file in self.output.rglob("*"):
            if output_file.is_file():
                self.assertNotIn(str(self.base), output_file.read_text(encoding="utf-8"))

    def test_dry_run_filters_and_writes_nothing(self):
        self.write("folder/a.pdf")
        self.write("folder/b.txt")
        self.write("other/c.txt")
        result = importer.run_import(
            self.source, self.output, mode="dry-run", extensions={"txt"},
            subfolder="folder", limit=1,
        )
        self.assertEqual(result["examined_count"], 1)
        self.assertFalse(self.output.exists())

    def test_balanced_limits_by_extension(self):
        for index in range(4):
            self.write(f"pdf/{index}.pdf", b"synthetic")
            self.write(f"text/{index}.txt", b"synthetic")
        result = importer.run_import(
            self.source, self.output, mode="dry-run",
            extension_limits={"pdf": 2, "txt": 3},
        )
        self.assertEqual(result["examined_count"], 5)


if __name__ == "__main__":
    unittest.main()
